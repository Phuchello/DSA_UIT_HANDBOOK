import os
import sys
import io
import zipfile
import pypdf
import pdfplumber
import pptx
import docx
import openpyxl
import json
import traceback

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')


RESEARCH_DIR = r"C:\Users\lyle3\.gemini\antigravity\scratch\IT003_DSA_BOOK\research"
EXTRACTED_DIR = os.path.join(RESEARCH_DIR, "extracted")
os.makedirs(EXTRACTED_DIR, exist_ok=True)

TARGET_DIRS = [
    r"C:\Users\lyle3\Downloads",
    r"C:\Users\lyle3\OneDrive\Documents\Món học\DSA"
]

inventory = []

def extract_pdf(pdf_path, out_txt_path):
    text_content = []
    try:
        reader = pypdf.PdfReader(pdf_path)
        num_pages = len(reader.pages)
        for i, page in enumerate(reader.pages):
            txt = page.extract_text() or ""
            text_content.append(f"--- PAGE {i+1} ---\n{txt}")
        full_text = "\n".join(text_content)
        with open(out_txt_path, "w", encoding="utf-8") as f:
            f.write(full_text)
        return num_pages, len(full_text)
    except Exception as e:
        print(f"Error extracting PDF {pdf_path}: {e}")
        return 0, 0

def extract_pptx(pptx_path, out_txt_path):
    text_content = []
    try:
        prs = pptx.Presentation(pptx_path)
        for i, slide in enumerate(prs.slides):
            slide_txt = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        slide_txt.append(paragraph.text)
            text_content.append(f"--- SLIDE {i+1} ---\n" + "\n".join(slide_txt))
        full_text = "\n".join(text_content)
        with open(out_txt_path, "w", encoding="utf-8") as f:
            f.write(full_text)
        return len(prs.slides), len(full_text)
    except Exception as e:
        print(f"Error extracting PPTX {pptx_path}: {e}")
        return 0, 0

def extract_docx(docx_path, out_txt_path):
    try:
        doc = docx.Document(docx_path)
        full_text = "\n".join([p.text for p in doc.paragraphs])
        with open(out_txt_path, "w", encoding="utf-8") as f:
            f.write(full_text)
        return len(doc.paragraphs), len(full_text)
    except Exception as e:
        print(f"Error extracting DOCX {docx_path}: {e}")
        return 0, 0

def extract_xlsx(xlsx_path, out_txt_path):
    try:
        wb = openpyxl.load_workbook(xlsx_path, data_only=True)
        sheets_text = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_rows = []
            for row in sheet.iter_rows(values_only=True):
                row_str = "\t".join([str(cell) if cell is not None else "" for cell in row])
                if row_str.strip():
                    sheet_rows.append(row_str)
            sheets_text.append(f"--- SHEET: {sheet_name} ---\n" + "\n".join(sheet_rows))
        full_text = "\n\n".join(sheets_text)
        with open(out_txt_path, "w", encoding="utf-8") as f:
            f.write(full_text)
        return len(wb.sheetnames), len(full_text)
    except Exception as e:
        print(f"Error extracting XLSX {xlsx_path}: {e}")
        return 0, 0

def process_all():
    extracted_count = 0
    unzip_dir = os.path.join(RESEARCH_DIR, "unzipped")
    os.makedirs(unzip_dir, exist_ok=True)
    
    for base_dir in TARGET_DIRS:
        if not os.path.exists(base_dir):
            continue
        for root, dirs, files in os.walk(base_dir):
            for file in files:
                filepath = os.path.join(root, file)
                ext = os.path.splitext(file)[1].lower()
                
                if ext == ".zip":
                    try:
                        with zipfile.ZipFile(filepath, 'r') as zip_ref:
                            sub_unzip = os.path.join(unzip_dir, os.path.splitext(file)[0])
                            zip_ref.extractall(sub_unzip)
                            print(f"Unzipped {file} to {sub_unzip}")
                    except Exception as e:
                        print(f"Failed to unzip {file}: {e}")

    all_scan_dirs = TARGET_DIRS + [unzip_dir]
    seen_files = set()

    for base_dir in all_scan_dirs:
        if not os.path.exists(base_dir):
            continue
        for root, dirs, files in os.walk(base_dir):
            for file in files:
                filepath = os.path.join(root, file)
                if filepath in seen_files:
                    continue
                seen_files.add(filepath)
                ext = os.path.splitext(file)[1].lower()
                
                if ext not in [".pdf", ".pptx", ".docx", ".xlsx", ".cpp", ".c", ".h", ".txt", ".md"]:
                    continue

                safe_name = "".join([c if c.isalnum() else "_" for c in os.path.splitext(file)[0]]) + ext.replace(".", "_") + ".txt"
                out_txt_path = os.path.join(EXTRACTED_DIR, safe_name)

                units = 0
                char_count = 0

                if ext == ".pdf":
                    units, char_count = extract_pdf(filepath, out_txt_path)
                elif ext == ".pptx":
                    units, char_count = extract_pptx(filepath, out_txt_path)
                elif ext == ".docx":
                    units, char_count = extract_docx(filepath, out_txt_path)
                elif ext == ".xlsx":
                    units, char_count = extract_xlsx(filepath, out_txt_path)
                elif ext in [".cpp", ".c", ".h", ".txt", ".md"]:
                    try:
                        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                            content = f.read()
                        with open(out_txt_path, "w", encoding="utf-8") as f:
                            f.write(content)
                        units = len(content.splitlines())
                        char_count = len(content)
                    except Exception as e:
                        print(f"Error copying txt/code {filepath}: {e}")

                if char_count > 0:
                    extracted_count += 1
                    inventory.append({
                        "file_name": file,
                        "file_path": filepath,
                        "ext": ext,
                        "extracted_txt": out_txt_path,
                        "units": units,
                        "char_count": char_count
                    })

    with open(os.path.join(RESEARCH_DIR, "inventory.json"), "w", encoding="utf-8") as f:
        json.dump(inventory, f, ensure_ascii=False, indent=2)

    print(f"Total files successfully extracted: {extracted_count}")

if __name__ == "__main__":
    process_all()
